from pptx import Presentation


prs = Presentation(r"C:\Users\Administrator\Desktop\moban.pptx")
print("slides", len(prs.slides), "layouts", len(prs.slide_layouts), "size", prs.slide_width, prs.slide_height)

for i, layout in enumerate(prs.slide_layouts):
    placeholders = []
    for ph in layout.placeholders:
        placeholders.append((ph.placeholder_format.idx, str(ph.placeholder_format.type), ph.name))
    print("LAYOUT", i, layout.name, placeholders)

for si, slide in enumerate(prs.slides):
    print("SLIDE", si)
    for sh in list(slide.shapes)[:20]:
        text = sh.text[:30].replace("\n", " ") if getattr(sh, "has_text_frame", False) else ""
        print("  ", sh.shape_type, sh.name, "text=", text)

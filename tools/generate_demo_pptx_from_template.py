from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = Path(r"C:\Users\Administrator\Desktop\moban.pptx")
OUT = ROOT / "8分钟演示文稿PPT.pptx"


PROJECT_TITLE = "基于特色创新项目的岭南红橙病虫害智能监测与校本教研一体化平台"


def duplicate_slide(prs: Presentation, source_index: int):
    """Duplicate a slide while preserving template shapes, pictures, and relationships."""
    source = prs.slides[source_index]
    blank_layout = prs.slide_layouts[6]
    dest = prs.slides.add_slide(blank_layout)

    for shape in source.shapes:
        dest.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)

    return dest


def remove_original_slides(prs: Presentation, count: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    for _ in range(count):
        r_id = sld_id_lst[0].rId
        prs.part.drop_rel(r_id)
        del sld_id_lst[0]


def text_shapes(slide):
    shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            continue
        if getattr(shape, "has_text_frame", False):
            shapes.append(shape)
    return shapes


def set_text(shape, text: str, size: int | None = None, bold: bool | None = None) -> None:
    shape.text_frame.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
        p.text = line
        for run in p.runs:
            run.font.name = "楷体"
            run.font.size = Pt(size or 24)
            if bold is not None:
                run.font.bold = bold


def fill_slide(slide, texts: list[str], sizes: list[int] | None = None, bolds: list[bool] | None = None) -> None:
    shapes = text_shapes(slide)
    for idx, text in enumerate(texts):
        if idx >= len(shapes):
            break
        set_text(
            shapes[idx],
            text,
            size=(sizes[idx] if sizes and idx < len(sizes) else None),
            bold=(bolds[idx] if bolds and idx < len(bolds) else None),
        )
    for idx in range(len(texts), len(shapes)):
        # Hide leftover template helper text without deleting decorative shapes.
        set_text(shapes[idx], "", size=1)


SLIDES = [
    {
        "template": 0,
        "texts": [
            PROJECT_TITLE,
            "广东文理职业学院",
            "演示视频：PPT + 录屏 + 解说",
        ],
        "sizes": [34, 24, 22],
        "bolds": [True, True, True],
    },
    {
        "template": 1,
        "texts": [
            "演示结构（总时长不超过8分钟）\n"
            "1. 案例概述：应用场景与主要问题（2分钟以内）\n"
            "2. 实现功能：系统功能与应用过程（5分钟以内）\n"
            "3. 应用情况：应用成效与推广价值（1分钟以内）"
        ],
        "sizes": [26],
        "bolds": [True],
    },
    {
        "template": 3,
        "texts": [
            "一、案例概述：背景与应用场景",
            "案例概述",
            "面向高职涉农物联网专业实训教学，依托广东省特色创新类项目（2025KQNCX312），将一线岭南红橙病虫害图像与物联网环境数据转化为校本教研资源。\n"
            "应用场景：课堂案例教学、实训室检测、田间图片分析、教研资料生成。"
        ],
        "sizes": [28, 22, 22],
        "bolds": [True, True, False],
    },
    {
        "template": 4,
        "texts": [
            "二、案例概述：解决的主要问题",
            "1. 科研成果转化慢，课题成果难以及时进入课堂\n"
            "2. 实训教学缺乏真实复杂农业生产情境\n"
            "3. 学生识别能力、实训过程和学习成果评价维度单一\n"
            "4. 教师备课、案例整理和台账归档耗时高",
            "案例概述",
        ],
        "sizes": [28, 23, 22],
        "bolds": [True, False, True],
    },
    {
        "template": 2,
        "texts": [
            "总体解决思路",
            "真实红橙病虫害图像\n物联网温湿度数据\nYOLOv8s-XH 本地识别\nSQLite 校本知识库推理\n三位一体防治方案\n教学案例与实训指导生成\n台账归档与评价",
            "监测",
            "识别",
            "决策",
            "教研",
        ],
        "sizes": [30, 21, 20, 20, 20, 20],
        "bolds": [True, False, True, True, True, True],
    },
    {
        "template": 3,
        "texts": [
            "三、实现功能：系统总体架构",
            "实现功能",
            "前端交互：PyQt5 桌面端，大字体、高对比度、单键触发\n"
            "视觉识别：YOLOv8s-XH 本地离线推理\n"
            "数据存储：SQLite 校本知识库与历史台账\n"
            "报表导出：Excel / PDF / Markdown\n"
            "教研生成：模板生成 + 可选大模型增强"
        ],
        "sizes": [28, 22, 22],
        "bolds": [True, True, False],
    },
    {
        "template": 4,
        "texts": [
            "四、实现功能：病虫害智能识别",
            "1. 支持图片、视频、摄像头检测\n"
            "2. 秒级识别12类高发病虫害\n"
            "3. 输出检测框、类别、置信度\n"
            "4. 支持红蜘蛛、蚜虫、木虱等群聚虫害计数\n"
            "5. 自动判断轻度、中度、重度",
            "实现功能",
        ],
        "sizes": [28, 23, 22],
        "bolds": [True, False, True],
    },
    {
        "template": 3,
        "texts": [
            "五、实现功能：绿色防治知识库推理",
            "实现功能",
            "联动病虫害类别、物候期、温湿度等多维数据，匹配 SQLite 校本知识库。\n"
            "一键生成物理、生物、合规化学“三位一体”防治方案，突出安全间隔期（PHI），对黄龙病等高危病害给出强提醒。"
        ],
        "sizes": [28, 22, 22],
        "bolds": [True, True, False],
    },
    {
        "template": 4,
        "texts": [
            "六、实现功能：农情台账与报表导出",
            "1. 自动保存检测记录，形成真实农情台账\n"
            "2. 支持按农户、果园、病虫害、时间筛选\n"
            "3. 3秒内离线生成 Excel 教学台账\n"
            "4. 支持 PDF 分析报告，服务课堂评价与教研归档",
            "实现功能",
        ],
        "sizes": [28, 23, 22],
        "bolds": [True, False, True],
    },
    {
        "template": 3,
        "texts": [
            "七、实现功能：校本教研内容生成",
            "实现功能",
            "基于最近一次检测结果，生成《岭南红橙实时防害教学案例》或《学生农技实训指导意见》。\n"
            "支持教师审核、修改、导出；可选接入 DeepSeek、豆包或本地大模型，将检测结果转化为课堂案例、实训任务和考核材料。"
        ],
        "sizes": [28, 22, 22],
        "bolds": [True, True, False],
    },
    {
        "template": 2,
        "texts": [
            "应用过程演示闭环",
            "导入红橙病虫害图片\nYOLOv8s-XH 本地识别\n选择物候期/读取环境数据\n知识库生成防治方案\n保存检测台账\n导出报表\n生成教学案例",
            "采集",
            "识别",
            "推理",
            "教研",
        ],
        "sizes": [30, 21, 20, 20, 20, 20],
        "bolds": [True, False, True, True, True, True],
    },
    {
        "template": 6,
        "texts": [
            "八、应用情况：成效与影响力",
            "平台已面向真实教学场景应用，沉淀红橙病虫害识别案例和农情台账。\n"
            "成效：提升教师案例备课效率，增强学生对真实农业场景的理解，支撑涉农物联网专业实训教学、过程评价和校本教研。",
            "应用情况",
            "科教融合：省特色创新项目转化为教学资源\n边缘智能：断网可运行，部署成本低\n教研闭环：真实数据自动生成实训指导\n推广价值：可拓展到柑橘、茶叶、蔬菜等农业实训场景",
        ],
        "sizes": [28, 23, 22, 21],
        "bolds": [True, False, True, False],
    },
    {
        "template": 8,
        "texts": [
            "谢谢！",
        ],
        "sizes": [44],
        "bolds": [True],
    },
]


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(TEMPLATE)
    prs = Presentation(str(TEMPLATE))
    original_count = len(prs.slides)

    new_slides = []
    for spec in SLIDES:
        new_slides.append(duplicate_slide(prs, spec["template"]))

    remove_original_slides(prs, original_count)

    for slide, spec in zip(new_slides, SLIDES):
        fill_slide(slide, spec["texts"], spec.get("sizes"), spec.get("bolds"))

    prs.save(str(OUT))
    print(f"generated: {OUT}")


if __name__ == "__main__":
    main()
